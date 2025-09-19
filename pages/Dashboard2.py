import streamlit as st
import pandas as pd
import plotly.express as px
import plotly.graph_objects as go
from config.thresholds import VACCINE_THRESHOLDS
from pptx import Presentation
from io import BytesIO

# --- Custom CSS for white background and black text ---
st.markdown("""
<style>
/* Overall page layout and styling - WHITE BACKGROUND */
.stApp {
    padding-top: 1rem;
    background: white !important;
    color: #000000 !important;
    min-height: 100vh;
}

/* Header styling - Professional with black text */
.main-header-container {
    background: white;
    padding: 1.5rem;
    border-radius: 12px;
    margin-bottom: 1.5rem;
    box-shadow: 0 6px 20px rgba(0, 0, 0, 0.08);
    text-align: center;
    border: 1px solid #e0e0e0;
}

.main-header-container h1 {
    color: #000000 !important;
    margin: 0;
    font-size: 1.8rem;
    font-weight: 700;
    letter-spacing: 0.5px;
}

.main-header-container p {
    color: #404040 !important;
    margin: 0.5rem 0 0 0;
    font-size: 1rem;
}

/* Enhanced metric cards - SMALLER */
.custom-metric-box {
    background: white;
    padding: 1rem;
    border-radius: 12px;
    text-align: center;
    margin-bottom: 1rem;
    box-shadow: 0 4px 12px rgba(0, 0, 0, 0.06);
    border: 1px solid #e0e0e0;
    transition: all 0.3s ease;
    position: relative;
    overflow: hidden;
    height: 120px;
    display: flex;
    flex-direction: column;
    justify-content: center;
}

.custom-metric-box::before {
    content: '';
    position: absolute;
    top: 0;
    left: 0;
    width: 3px;
    height: 100%;
    background: linear-gradient(135deg, #3498db 0%, #2c3e50 100%);
}

.custom-metric-box:hover {
    transform: translateY(-2px);
    box-shadow: 0 6px 16px rgba(0, 0, 0, 0.08);
}

.custom-metric-label {
    font-size: 0.75rem;
    font-weight: 600;
    color: #404040 !important;
    margin-bottom: 0.5rem;
    text-transform: uppercase;
    letter-spacing: 0.5px;
}

.custom-metric-value {
    font-size: 1.5rem;
    font-weight: 700;
    color: #000000 !important;
}

/* Section headers */
.stSubheader {
    color: #000000 !important;
    font-weight: 700 !important;
    border-bottom: 2px solid #3498db;
    padding-bottom: 0.5rem;
    margin-top: 2rem !important;
    font-size: 1.4rem !important;
}

/* Chart titles and labels - BLACK TEXT */
.js-plotly-plot .plotly .gtitle {
    color: #000000 !important;
    font-weight: 700 !important;
    font-size: 1.1rem !important;
    text-align: center;
}

.js-plotly-plot .plotly .xtitle, .js-plotly-plot .plotly .ytitle {
    color: #000000 !important;
    font-weight: 600 !important;
    font-size: 0.9rem !important;
}

.js-plotly-plot .plotly .legend text {
    color: #000000 !important;
    font-weight: 500 !important;
    font-size: 0.8rem !important;
}

/* Enhanced chart containers */
.js-plotly-plot {
    border-radius: 12px;
    box-shadow: 0 4px 12px rgba(0, 0, 0, 0.06);
    background: white;
    padding: 1rem;
    border: 1px solid #e0e0e0;
}

/* Warning and info message styling - FIXED VISIBILITY */
.stWarning {
    background: #fff3cd !important;
    border: 1px solid #ffd166 !important;
    border-radius: 10px !important;
    color: #856404 !important;
    padding: 1rem !important;
    box-shadow: 0 3px 10px rgba(255, 193, 7, 0.2) !important;
    font-size: 0.9rem !important;
}

.stInfo {
    background: #e9ecef !important;
    border: 1px solid #ced4da !important;
    border-radius: 10px !important;
    color: #000000 !important;
    padding: 1rem !important;
    box-shadow: 0 3px 10px rgba(0, 0, 0, 0.08) !important;
    font-size: 0.9rem !important;
}

.stSuccess {
    background: #d4edda !important;
    border: 1px solid #c3e6cb !important;
    border-radius: 10px !important;
    color: #155724 !important;
    padding: 1rem !important;
    box-shadow: 0 3px 10px rgba(40, 167, 69, 0.2) !important;
    font-size: 0.9rem !important;
}

.stWarning p, .stInfo p, .stSuccess p {
    color: inherit !important;
    font-weight: 500;
    margin: 0;
    font-size: 0.9rem;
}

.stWarning code, .stInfo code, .stSuccess code {
    color: inherit !important;
    background-color: rgba(255, 255, 255, 0.2) !important;
}

/* Enhanced sidebar */
.stSidebar {
    background: white !important;
    border-right: 1px solid #e0e0e0;
    box-shadow: 3px 0 15px rgba(0, 0, 0, 0.05);
}

.stSidebar .stSelectbox, .stSidebar .stSlider, .stSidebar .stTextInput {
    background: white !important;
    border-radius: 10px;
    padding: 0.6rem;
    border: 1px solid #e0e0e0;
    box-shadow: 0 2px 6px rgba(0, 0, 0, 0.04);
    color: #000000 !important;
    font-size: 0.9rem;
}

.stSidebar .stSelectbox label, .stSidebar .stSlider label, .stSidebar .stTextInput label {
    color: #000000 !important;
    font-weight: 600;
    font-size: 0.85rem;
}

.stSidebar .stHeader {
    color: #000000 !important;
    font-weight: 700;
    border-bottom: 2px solid #3498db;
    padding-bottom: 0.5rem;
    font-size: 1.1rem;
}

/* Enhanced buttons */
.stDownloadButton button {
    background: #3498db !important;
    color: white !important;
    border: none !important;
    border-radius: 10px !important;
    padding: 0.6rem 1.2rem !important;
    font-weight: 600 !important;
    font-size: 0.9rem;
    box-shadow: 0 3px 8px rgba(52, 152, 219, 0.3);
    transition: all 0.3s ease;
}

.stDownloadButton button:hover {
    transform: translateY(-1px);
    box-shadow: 0 5px 12px rgba(52, 152, 219, 0.4);
}

/* Enhanced dataframe */
.dataframe {
    border-radius: 10px;
    box-shadow: 0 3px 10px rgba(0, 0, 0, 0.04);
    border: 1px solid #e0e0e0;
    font-size: 0.85rem;
}

/* Enhanced expander */
.streamlit-expanderHeader {
    background: #f8f9fa;
    border-radius: 10px;
    padding: 0.8rem;
    font-weight: 600;
    color: #000000;
    border: 1px solid #e0e0e0;
    font-size: 0.9rem;
}

.streamlit-expanderContent {
    background: white;
    border-radius: 0 0 10px 10px;
    padding: 1.2rem;
    border: 1px solid #e0e0e0;
    border-top: none;
    font-size: 0.85rem;
}

/* Divider styling */
.stMarkdown hr {
    margin: 2rem 0;
    border: none;
    height: 2px;
    background: linear-gradient(90deg, transparent, #3498db, transparent);
}

/* Filter section styling - ENHANCED */
.filter-section {
    background: white !important;
    padding: 1.2rem;
    border-radius: 12px;
    margin-bottom: 1.2rem;
    border: 1px solid #e0e0e0;
    box-shadow: 0 3px 10px rgba(0, 0, 0, 0.04);
    color: #000000 !important;
}

.filter-section h2, .filter-section h3, .filter-section h4 {
    color: #000000 !important;
    font-weight: 700;
    font-size: 1rem;
}

.filter-section label {
    color: #000000 !important;
    font-weight: 600;
    font-size: 0.85rem;
}

/* Unmatched records styling - ENHANCED */
.unmatched-section {
    background: white;
    padding: 1.2rem;
    border-radius: 12px;
    border: 1px solid #e0e0e0;
    box-shadow: 0 3px 10px rgba(0, 0, 0, 0.04);
    color: #000000 !important;
    margin-bottom: 1rem;
}

.unmatched-section h3, .unmatched-section h4 {
    color: #000000 !important;
    font-weight: 700;
    border-bottom: 2px solid #3498db;
    padding-bottom: 0.5rem;
    margin-bottom: 1rem;
    font-size: 1rem;
}

/* Section header styling */
.section-header {
    background: white;
    padding: 1.2rem;
    border-radius: 12px;
    margin: 1.5rem 0;
    color: #000000 !important;
    text-align: center;
    border: 1px solid #e0e0e0;
    box-shadow: 0 4px 12px rgba(0, 0, 0, 0.06);
}

.section-header h3 {
    color: #000000 !important;
    margin: 0;
    font-weight: 700;
    font-size: 1.3rem;
}

/* Responsive adjustments */
@media (max-width: 768px) {
    .main-header-container h1 {
        font-size: 1.5rem;
    }
    
    .custom-metric-value {
        font-size: 1.3rem;
    }
    
    .stSubheader {
        font-size: 1.2rem !important;
    }
    
    .custom-metric-box {
        padding: 0.8rem;
        height: 110px;
    }
}

/* Animation for metrics */
@keyframes fadeInUp {
    from {
        opacity: 0;
        transform: translateY(15px);
    }
    to {
        opacity: 1;
        transform: translateY(0);
    }
}

.custom-metric-box {
    animation: fadeInUp 0.5s ease-out;
}

/* Grid line enhancements - REMOVED gridlines */
.js-plotly-plot .plotly .gridlayer .grid {
    stroke: none !important;
}

.js-plotly-plot .plotly .gridlayer .xgrid, 
.js-plotly-plot .plotly .gridlayer .ygrid {
    stroke: none !important;
}

/* Custom scrollbar */
::-webkit-scrollbar {
    width: 6px;
}

::-webkit-scrollbar-track {
    background: #f1f1f1;
    border-radius: 3px;
}

::-webkit-scrollbar-thumb {
    background: #3498db;
    border-radius: 3px;
}

::-webkit-scrollbar-thumb:hover {
    background: #2980b9;
}

/* Fix for all text elements */
.stSidebar * {
    color: #000000 !important;
}

.stSidebar .st-bb {
    background-color: white !important;
}

.stSidebar .st-at {
    background-color: white !important;
}

/* Chart specific improvements */
.plotly .modebar {
    background: rgba(255, 255, 255, 0.9) !important;
    border: 1px solid #e0e0e0 !important;
    border-radius: 6px !important;
}

.plotly .modebar-btn {
    color: #000000 !important;
}

.plotly .hovertext {
    background: white !important;
    border: 1px solid #e0e0e0 !important;
    color: #000000 !important;
    font-size: 0.8rem !important;
    border-radius: 6px !important;
    box-shadow: 0 3px 10px rgba(0, 0, 0, 0.1) !important;
}
</style>
""", unsafe_allow_html=True)

st.set_page_config(
    page_title="Vaccine Utilization Analytics Dashboard",
    layout="wide",
    page_icon="💉"
)

# --- Helper Functions ---
def to_excel(df):
    """Convert DataFrame to Excel file for download"""
    output = BytesIO()
    with pd.ExcelWriter(output, engine='xlsxwriter') as writer:
        df.to_excel(writer, index=False, sheet_name='Sheet1')
    processed_data = output.getvalue()
    return processed_data

def count_extremity(df, vaccine):
    """Calculates high and low extremity counts based on user-defined thresholds."""
    rate_col = f"{vaccine}_Utilization_Rate"
    if rate_col not in df.columns:
        return None
    
    # User-provided thresholds (as decimals)
    thresholds = {
        "BCG": {"high": 1.20, "low": 0.30},
        "IPV": {"high": 1.30, "low": 0.75},
        "Measles": {"high": 1.25, "low": 0.50},
        "Penta": {"high": 1.30, "low": 0.80},
        "Rota": {"high": 1.30, "low": 0.75},
    }
    
    if vaccine not in thresholds:
        return None

    high_count = len(df[df[rate_col] / 100 > thresholds[vaccine]["high"]])
    low_count = len(df[df[rate_col] / 100 < thresholds[vaccine]["low"]])
    
    return high_count, low_count

def setup_filters(df_all):
    """Set up sidebar filters and return filtered DataFrame"""
    st.sidebar.markdown('<div class="filter-section">', unsafe_allow_html=True)
    st.sidebar.header("🧪 Filter Data")
    
    # Region filter
    regions = sorted(df_all["Region_Admin"].unique())
    selected_region = st.sidebar.selectbox("Select Region", ["All"] + regions)
    
    # Zone filter (chained to Region)
    filtered_zones = df_all[df_all["Region_Admin"] == selected_region]["Zone_Admin"].unique() if selected_region != "All" else df_all["Zone_Admin"].unique()
    zones = sorted(filtered_zones)
    selected_zone = st.sidebar.selectbox("Select Zone", ["All"] + zones)
    
    # Woreda filter (chained to Zone)
    if selected_zone != "All":
        filtered_woredas = df_all[(df_all["Region_Admin"] == selected_region) & (df_all["Zone_Admin"] == selected_zone)]["Woreda_Admin"].unique()
    elif selected_region != "All":
        filtered_woredas = df_all[df_all["Region_Admin"] == selected_region]["Woreda_Admin"].unique()
    else:
        filtered_woredas = df_all["Woreda_Admin"].unique()
    woredas = sorted(filtered_woredas)
    selected_woreda = st.sidebar.selectbox("Select Woreda", ["All"] + woredas)

    # Period filter
    periods = sorted(df_all["Period"].unique())
    selected_period = st.sidebar.selectbox("Select Period", ["All"] + periods)
    
    # Vaccine filter
    vaccines = ["BCG", "IPV", "Measles", "Penta", "Rota"]
    selected_vaccine = st.sidebar.selectbox("Select Vaccine", ["All"] + vaccines)
    st.sidebar.markdown('</div>', unsafe_allow_html=True)
    
    # Filtering Logic
    filtered_df = df_all
    if selected_region != "All":
        filtered_df = filtered_df[filtered_df["Region_Admin"] == selected_region]
    if selected_zone != "All":
        filtered_df = filtered_df[filtered_df["Zone_Admin"] == selected_zone]
    if selected_woreda != "All":
        filtered_df = filtered_df[filtered_df["Woreda_Admin"] == selected_woreda]
    if selected_period != "All":
        filtered_df = filtered_df[filtered_df["Period"] == selected_period]
        
    return filtered_df, selected_vaccine, vaccines

def display_kpis(filtered_df, selected_vaccine, vaccines):
    """Display Key Performance Indicators"""
    total_woredas = filtered_df["Woreda_Admin"].nunique()
    
    if selected_vaccine != "All":
        dist_col = f"{selected_vaccine}_Distributed"
        admin_col = f"{selected_vaccine}_Administered"
        if dist_col in filtered_df.columns and admin_col in filtered_df.columns:
            total_admin = filtered_df[admin_col].sum()
            total_dist = filtered_df[dist_col].sum()
        else:
            st.warning(f"Data for '{selected_vaccine}' is not available in the processed files.")
            return
    else:
        dist_cols = [f"{v}_Distributed" for v in vaccines]
        admin_cols = [f"{v}_Administered" for v in vaccines]
        
        existing_dist_cols = [col for col in dist_cols if col in filtered_df.columns]
        existing_admin_cols = [col for col in admin_cols if col in filtered_df.columns]
        
        total_dist = filtered_df[existing_dist_cols].sum().sum()
        total_admin = filtered_df[existing_admin_cols].sum().sum()
    
    utilization_rate = (total_admin / total_dist) * 100 if total_dist > 0 else 0
    
    st.markdown("---")
    st.markdown('<div class="section-header"><h3>📊 Performance Overview</h3></div>', unsafe_allow_html=True)
    
    col1, col2, col3, col4 = st.columns(4)
    with col1:
        st.markdown(f'<div class="custom-metric-box"><div class="custom-metric-label">Total Woredas</div><div class="custom-metric-value">{total_woredas}</div></div>', unsafe_allow_html=True)
    with col2:
        st.markdown(f'<div class="custom-metric-box"><div class="custom-metric-label">Total Doses Distributed</div><div class="custom-metric-value">{total_dist:,.0f}</div></div>', unsafe_allow_html=True)
    with col3:
        st.markdown(f'<div class="custom-metric-box"><div class="custom-metric-label">Total Doses Administered</div><div class="custom-metric-value">{total_admin:,.0f}</div></div>', unsafe_allow_html=True)
    with col4:
        st.markdown(f'<div class="custom-metric-box"><div class="custom-metric-label">Utilization Rate</div><div class="custom-metric-value">{utilization_rate:.2f}%</div></div>', unsafe_allow_html=True)
    st.markdown("---")

def display_extremities(filtered_df, vaccines):
    """Display extremity counts for each vaccine"""
    st.markdown('<div class="section-header"><h3>🚨 Extreme Utilization Rates</h3></div>', unsafe_allow_html=True)
    
    counts_col1, counts_col2, counts_col3, counts_col4, counts_col5 = st.columns(5)
    
    for i, vaccine in enumerate(vaccines):
        counts = count_extremity(filtered_df, vaccine)
        if counts:
            with [counts_col1, counts_col2, counts_col3, counts_col4, counts_col5][i]:
                st.markdown(f'<div class="custom-metric-box"><div class="custom-metric-label">{vaccine}</div><div class="custom-metric-value">{counts[0]}↑ | {counts[1]}↓</div></div>', unsafe_allow_html=True)
    st.markdown("---")

def display_extreme_utilization_table(filtered_df, selected_vaccine):
    """Display table of extreme utilization by region and zone"""
    st.markdown('<div class="section-header"><h3>📋 Extreme Utilization by Region and Zone</h3></div>', unsafe_allow_html=True)
    
    if selected_vaccine == "All":
        st.info("Please select a specific vaccine to view this table.")
        return
    
    # Check if the required column exists
    rate_col = f"{selected_vaccine}_Utilization_Rate"
    if rate_col not in filtered_df.columns:
        st.warning(f"Utilization data for {selected_vaccine} is not available in the processed files.")
        return
    
    # User-provided thresholds (as decimals)
    thresholds = {
        "BCG": {"high": 1.20, "low": 0.30},
        "IPV": {"high": 1.30, "low": 0.75},
        "Measles": {"high": 1.25, "low": 0.50},
        "Penta": {"high": 1.30, "low": 0.80},
        "Rota": {"high": 1.30, "low": 0.75},
    }

    # Create a copy of the filtered data
    extreme_df = filtered_df.copy()
    
    # Determine grouping columns based on filter selection
    if "All" in st.session_state.get('selected_region', 'All'):
        group_cols = ["Region_Admin"]
    else:
        group_cols = ["Region_Admin", "Zone_Admin"]
    
    # Calculate counts for each category
    extreme_df["High_Extremity"] = extreme_df[rate_col] / 100 > thresholds[selected_vaccine]["high"]
    extreme_df["Low_Extremity"] = extreme_df[rate_col] / 100 < thresholds[selected_vaccine]["low"]

    extreme_summary = extreme_df.groupby(group_cols).agg(
        total_woredas=("Woreda_Admin", "nunique"),
        high_extremity_count=("High_Extremity", "sum"),
        low_extremity_count=("Low_Extremity", "sum")
    ).reset_index()
    
    # Rename columns for display
    extreme_summary.rename(columns={
        "Region_Admin": "Region",
        "Zone_Admin": "Zone",
        "total_woredas": "Total Woredas",
        "high_extremity_count": "High Extremity",
        "low_extremity_count": "Low Extremity",
    }, inplace=True)

    st.dataframe(extreme_summary, use_container_width=True, hide_index=True)
    
    st.download_button(
        label="📥 Download Data as Excel",
        data=to_excel(extreme_summary),
        file_name=f"Extreme_Utilization_Data_{selected_vaccine}.xlsx",
        mime="application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
        use_container_width=True
    )

def display_utilization_chart(filtered_df, selected_vaccine):
    """Display utilization rate chart by Woreda"""
    st.markdown('<div class="section-header"><h3>📈 Utilization Rate by Woreda</h3></div>', unsafe_allow_html=True)
    
    if selected_vaccine == "All":
        st.info("Please select a specific vaccine to view utilization rate plots.")
        return
    
    # Check if the required column exists
    rate_col = f"{selected_vaccine}_Utilization_Rate"
    if rate_col not in filtered_df.columns:
        st.warning(f"Utilization data for {selected_vaccine} is not available in the processed files.")
        return
    
    # Create a scatter plot of utilization rates
    fig = px.scatter(
        filtered_df,
        x="Woreda_Admin",
        y=rate_col,
        title=f"{selected_vaccine} Utilization Rate by Woreda",
        labels={"Woreda_Admin": "Woreda", rate_col: "Utilization Rate (%)"},
        color_discrete_sequence=["#3498db"]
    )
    
    # Add threshold lines
    thresholds = {
        "BCG": {"high": 120, "low": 30},
        "IPV": {"high": 130, "low": 75},
        "Measles": {"high": 125, "low": 50},
        "Penta": {"high": 130, "low": 80},
        "Rota": {"high": 130, "low": 75},
    }
    
    if selected_vaccine in thresholds:
        fig.add_hline(y=thresholds[selected_vaccine]["high"], line_dash="dash", line_color="red", 
                     annotation_text="High Threshold", annotation_position="bottom right")
        fig.add_hline(y=thresholds[selected_vaccine]["low"], line_dash="dash", line_color="orange", 
                     annotation_text="Low Threshold", annotation_position="top right")
    
    # Improve visibility of labels and legends
    fig.update_layout(
        xaxis_title="Woreda", 
        yaxis_title="Utilization Rate (%)", 
        xaxis_tickangle=45,
        plot_bgcolor='#ffffff',
        paper_bgcolor='#ffffff',
        font=dict(color='#000000', size=12),
        legend=dict(
            orientation="h",
            yanchor="bottom",
            y=1.02,
            xanchor="right",
            x=1,
            font=dict(size=11, color='#000000'),
            bgcolor='rgba(255,255,255,0.9)',
            bordercolor='#e0e0e0',
            borderwidth=1
        ),
        hoverlabel=dict(
            bgcolor="white",
            font_size=11,
            font_family="Arial",
            font_color="#000000",
            bordercolor="#e0e0e0"
        )
    )
    
    # Improve x-axis label visibility
    fig.update_xaxes(tickfont=dict(size=9, color='#000000'), showgrid=False)
    fig.update_yaxes(tickfont=dict(color='#000000'), showgrid=False)
    
    st.plotly_chart(fig, use_container_width=True)

def display_unmatched_records():
    """Display unmatched records"""
    st.markdown('<div class="section-header"><h3>🔍 Unmatched Records</h3></div>', unsafe_allow_html=True)
    
    unmatched_col1, unmatched_col2 = st.columns(2)
    
    with unmatched_col1:
        st.markdown('<div class="unmatched-section">', unsafe_allow_html=True)
        if "unmatched_admin_df" in st.session_state and not st.session_state["unmatched_admin_df"].empty:
            st.write("**Unmatched Administered Records**")
            st.dataframe(st.session_state["unmatched_admin_df"], use_container_width=True)
        else:
            st.info("No unmatched administered records found.")
        st.markdown('</div>', unsafe_allow_html=True)
    
    with unmatched_col2:
        st.markdown('<div class="unmatched-section">', unsafe_allow_html=True)
        if "unmatched_dist_df" in st.session_state and not st.session_state["unmatched_dist_df"].empty:
            st.write("**Unmatched Distributed Records**")
            st.dataframe(st.session_state["unmatched_dist_df"], use_container_width=True)
        else:
            st.info("No unmatched distributed records found.")
        st.markdown('</div>', unsafe_allow_html=True)
            
    st.markdown("---")

def create_ppt(filtered_df, selected_vaccine):
    """Create PowerPoint report"""
    prs = Presentation()
    
    # Slide 1: Title Slide
    title_slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_slide_layout)
    title = slide.shapes.title
    title.text = "Immunization Triangulation Report"
    
    body = slide.placeholders[1]
    body.text = f"Report generated for {selected_vaccine}."

    # Slide 2: Recommendations for High Utilization
    rec_high_layout = prs.slide_layouts[1]
    rec_high_slide = prs.slides.add_slide(rec_high_layout)
    rec_high_slide.shapes.title.text = "Recommendations for High Utilization"
    
    rec_high_text = rec_high_slide.placeholders[1].text_frame
    
    p_high_bullet1 = rec_high_text.add_paragraph()
    p_high_bullet1.text = "Conduct a targeted audit of administered dose records for potential data entry errors."
    p_high_bullet1.level = 0

    p_high_bullet2 = rec_high_text.add_paragraph()
    p_high_bullet2.text = "Investigate potential lags in reporting of distributed doses."
    p_high_bullet2.level = 0
    
    p_high_bullet3 = rec_high_text.add_paragraph()
    p_high_bullet3.text = "Recommend a physical stock count at facilities to reconcile administered and distributed doses."
    p_high_bullet3.level = 0

    # Slide 3: Recommendations for Low Utilization
    rec_low_layout = prs.slide_layouts[1]
    rec_low_slide = prs.slides.add_slide(rec_low_layout)
    rec_low_slide.shapes.title.text = "Recommendations for Low Utilization"
    
    rec_low_text = rec_low_slide.placeholders[1].text_frame

    p_low_bullet1 = rec_low_text.add_paragraph()
    p_low_bullet1.text = "Assess inventory levels to prevent vaccine expiration due to overstocking."
    p_low_bullet1.level = 0

    p_low_bullet2 = rec_low_text.add_paragraph()
    p_low_bullet2.text = "Investigate if there are service delivery issues affecting vaccine demand."
    p_low_bullet2.level = 0

    p_low_bullet3 = rec_low_text.add_paragraph()
    p_low_bullet3.text = "Verify that administered doses are being reported accurately and on time."
    p_low_bullet3.level = 0

    # Slide 4: Recommendations for High Discrepancies
    rec_disc_layout = prs.slide_layouts[1]
    rec_disc_slide = prs.slides.add_slide(rec_disc_layout)
    rec_disc_slide.shapes.title.text = "Recommendations for High Discrepancies"
    
    rec_disc_text = rec_disc_slide.placeholders[1].text_frame
    
    p_disc_bullet1 = rec_disc_text.add_paragraph()
    p_disc_bullet1.text = "Provide training on accurate reporting for both administered and distributed doses."
    p_disc_bullet1.level = 0

    p_disc_bullet2 = rec_disc_text.add_paragraph()
    p_disc_bullet2.text = "Implement a process to regularly cross-reference data to catch discrepancies early."
    p_disc_bullet2.level = 0

    p_disc_bullet3 = rec_disc_text.add_paragraph()
    p_disc_bullet3.text = "Establish a feedback loop where Woredas are alerted to discrepancies."
    p_disc_bullet3.level = 0
    
    ppt_file = BytesIO()
    prs.save(ppt_file)
    ppt_file.seek(0)
    return ppt_file

def display_downloads(filtered_df, selected_vaccine):
    """Display download options"""
    ppt_buffer = create_ppt(filtered_df, selected_vaccine)
    st.download_button(
        label="📊 Download Report as PPT",
        data=ppt_buffer,
        file_name="immunization_report.pptx",
        mime="application/vnd.openxmlformats-officedocument.presentationml.presentation",
        use_container_width=True
    )

# --- Main Dashboard Logic ---
def main():
    """Main dashboard function"""
    # --- Header with Title Only ---
    st.markdown("""
    <div class="main-header-container">
        <h1>📊 Vaccine Discrepancies and Utilization Analysis</h1>
        <p>Comprehensive Vaccine Performance Monitoring</p>
    </div>
    """, unsafe_allow_html=True)

    if not st.session_state.get("authenticated", False):
        st.warning("Please log in on the main page to view this dashboard.")
        return

    if st.session_state.get("matched_df") is None:
        st.info("Please upload and process data on the Data Upload page to view this dashboard.")
        return
        
    df_all = st.session_state["matched_df"]
    
    # Check for required columns
    required_cols = ["Region_Admin", "Zone_Admin", "Woreda_Admin", "Period"]
    if not all(col in df_all.columns for col in required_cols):
        st.error("Essential columns (Region, Zone, Woreda, Period) are missing from the processed data.")
        return
    
    # Prepare data for calculation
    vaccines = ["BCG", "IPV", "Measles", "Penta", "Rota"]
    for vaccine in vaccines:
        admin_col = f"{vaccine}_Administered"
        dist_col = f"{vaccine}_Distributed"
        rate_col = f"{vaccine}_Utilization_Rate"
        if admin_col in df_all.columns and dist_col in df_all.columns:
            df_all[rate_col] = df_all[admin_col] / df_all[dist_col].replace(0, 1) * 100
            df_all[rate_col] = df_all[rate_col].clip(0, 1000)

    # Setup filters
    filtered_df, selected_vaccine, vaccines = setup_filters(df_all)

    if filtered_df.empty:
        st.warning("⚠️ No data found for the selected filters.")
        return

    # Display all dashboard components
    display_kpis(filtered_df, selected_vaccine, vaccines)
    display_extremities(filtered_df, vaccines)
    display_extreme_utilization_table(filtered_df, selected_vaccine)
    display_utilization_chart(filtered_df, selected_vaccine)
    display_unmatched_records()
    display_downloads(filtered_df, selected_vaccine)

# Run the main function
if __name__ == "__main__":
    main()
